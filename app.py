from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import os
import copy

app = Flask(__name__)

# --- Helper: Hybrid Slide Duplication ---
def duplicate_slide(prs, source_slide):
    """
    Duplicates a slide using a hybrid method:
    - Images: Re-inserted to establish correct relationships (fixes "Picture can't be displayed").
    - Shapes/Text: XML-cloned to preserve complex formatting and rotation.
    """
    # 1. Create a new slide using the SAME layout
    layout = source_slide.slide_layout
    dest = prs.slides.add_slide(layout)

    # 2. Clear default placeholders on the new slide
    for shape in list(dest.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # 3. Copy Background (if overridden on source)
    try:
        if source_slide.background.element is not None:
             dest.background.element = copy.deepcopy(source_slide.background.element)
    except:
        pass

    # 4. Loop through all shapes in the source slide
    for shape in source_slide.shapes:
        
        # --- CASE A: The Shape is an IMAGE (Type 13) ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # We must extract the image data and re-add it to create a valid relationship
            try:
                blob = shape.image.blob
                image_stream = io.BytesIO(blob)
                
                # Add picture at the exact same position and size
                new_pic = dest.shapes.add_picture(
                    image_stream, 
                    shape.left, shape.top, 
                    shape.width, shape.height
                )
                
                # Apply rotation if the original had it
                if hasattr(shape, 'rotation'):
                    new_pic.rotation = shape.rotation
                    
            except Exception as e:
                # Fallback: If image extraction fails, try XML copy (better than nothing)
                print(f"Image copy failed, trying XML fallback: {e}")
                new_element = copy.deepcopy(shape.element)
                dest.shapes._spTree.insert_element_before(new_element, 'p:extLst')

        # --- CASE B: The Shape is Text, Line, or AutoShape ---
        else:
            # Deep copy the XML. This preserves fonts, colors, borders, and rotation perfectly.
            new_element = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    return dest

# --- Main Processing Logic ---
def process_pptx_template(template_file, running_numbers, placeholder_text, items_per_slide):
    prs = Presentation(template_file)
    
    # Clean data
    codes = [line.strip() for line in running_numbers.strip().splitlines() if line.strip()]
    if not codes:
        return None

    # Chunk data
    chunks = [codes[i:i + items_per_slide] for i in range(0, len(codes), items_per_slide)]
    
    # --- PHASE 1: DUPLICATE SLIDES ---
    source_slide = prs.slides[0]
    needed_duplicates = len(chunks) - 1
    
    for _ in range(needed_duplicates):
        duplicate_slide(prs, source_slide)
        
    # --- PHASE 2: REPLACE DATA ---
    for i, chunk in enumerate(chunks):
        slide = prs.slides[i]
        chunk_index = 0
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder_text in run.text:
                        if chunk_index < len(chunk):
                            run.text = run.text.replace(placeholder_text, chunk[chunk_index], 1)
                            chunk_index += 1
                        else:
                            run.text = run.text.replace(placeholder_text, "")

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# --- Routes ---

@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    try:
        if 'template_file' not in request.files:
            return "No file uploaded", 400
            
        file = request.files['template_file']
        running_numbers = request.form.get('running_numbers', '')
        placeholder_text = request.form.get('placeholder_text', '{{NUM}}')
        items_per_slide = int(request.form.get('items_per_slide', 1))

        if file.filename == '':
            return "No selected file", 400

        output_pptx = process_pptx_template(file, running_numbers, placeholder_text, items_per_slide)

        if output_pptx:
            return send_file(
                output_pptx,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=f'Filled_{file.filename}'
            )
        else:
            return "Error: No running numbers provided.", 400

    except Exception as e:
        print(f"Error: {e}")
        return f"An error occurred: {str(e)}", 500

if __name__ == '__main__':
    if not os.path.exists('templates'):
        os.makedirs('templates')
    app.run(debug=True)